from pptx import Presentation
from pptx.util import Pt

# Crear la presentación
prs = Presentation()

# Tamaños estándar
title_font_size = Pt(40)
content_font_size = Pt(24)

# Función para agregar diapositivas
def add_slide(title, content_lines):
    slide_layout = prs.slide_layouts[1]  # Título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = "\n".join(content_lines)

    for paragraph in content_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = content_font_size

# Diapositiva 1: Portada
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Modelo de IA DeepSeek"
slide.placeholders[1].text = "Aplicación para la clase de Compiladores - IPN\nAgosto 2025"

# Diapositiva 2: ¿Qué es DeepSeek?
add_slide("¿Qué es DeepSeek?", [
    "DeepSeek es una familia de modelos de lenguaje desarrollada por DeepSeek AI.",
    "Incluye variantes como DeepSeek LLM y DeepSeek Coder.",
    "Optimizado para tareas generales de lenguaje y tareas de programación.",
    "Compite con modelos como GPT-4, Claude y Gemini."
])

# Diapositiva 3: Arquitectura de DeepSeek
add_slide("Arquitectura de DeepSeek", [
    "Basado en la arquitectura Transformer, al igual que GPT.",
    "Entrenado con trillones de tokens, tanto de lenguaje natural como de código.",
    "Versión DeepSeek Coder: especializada en comprensión y generación de código.",
    "Usa técnicas de preentrenamiento y fine-tuning específicas para programación."
])

# Diapositiva 4: Características clave
add_slide("Características clave", [
    "DeepSeek Coder soporta más de 80 lenguajes de programación.",
    "Entrenamiento mixto: 87% código + 13% lenguaje natural.",
    "Alto rendimiento en benchmarks como HumanEval y MBPP.",
    "Generación estructurada, autocompletado de código, y comprensión semántica."
])

# Diapositiva 5: Comparación con otros modelos
add_slide("Comparación con otros modelos", [
    "DeepSeek Coder es comparable a modelos como GPT-4 Code Interpreter.",
    "Muestra mejor desempeño en generación de código específico y tareas técnicas.",
    "Open-source (disponible en HuggingFace), lo cual facilita su uso académico.",
    "Menor latencia en ciertas tareas frente a modelos comerciales cerrados."
])

# Diapositiva 6: Aplicaciones en Compiladores
add_slide("Aplicaciones en Compiladores", [
    "Generación de código a partir de descripciones en lenguaje natural.",
    "Detección de errores sintácticos y sugerencia de correcciones.",
    "Análisis semántico automatizado para entender intenciones del código.",
    "Traducción entre lenguajes de programación (transcompilación)."
])

# Diapositiva 7: Conclusión
add_slide("Conclusión", [
    "DeepSeek es un modelo potente, especialmente útil en programación.",
    "Su arquitectura y enfoque lo hacen ideal para apoyar cursos de compiladores.",
    "Provee herramientas para analizar, entender y generar código de forma eficaz.",
    "El acceso abierto permite su integración en proyectos educativos y de investigación."
])

# Guardar el archivo
prs.save("Modelo_IA_DeepSeek_Compiladores_IPN.pptx")
